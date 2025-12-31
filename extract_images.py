import pptx
import sys
import os

pptx_path = sys.argv[1]
output_dir = "extracted_slides"
os.makedirs(output_dir, exist_ok=True)

prs = pptx.Presentation(pptx_path)

for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13: # Picture
            image = shape.image
            # ensure filename sorts correctly
            filename = f"slide_{i+1:03d}.{image.ext}" 
            with open(os.path.join(output_dir, filename), "wb") as f:
                f.write(image.blob)
