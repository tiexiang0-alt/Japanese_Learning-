import pptx
import sys
import os

try:
    prs = pptx.Presentation(sys.argv[1])
    text_found = False
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(f"  Text: {text}")
                    text_found = True
            if shape.shape_type == 13: # Picture
                 print(f"  Found Picture")
    
    if not text_found:
        print("No text found. Likely flattened images.")

except Exception as e:
    print(f"Error: {e}")
